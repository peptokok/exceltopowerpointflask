from flask import Flask, render_template, request, jsonify, send_file
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import base64
import io
import os

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload-excel', methods=['POST'])
def upload_excel():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith(('.xlsx', '.xls')):
        return jsonify({'error': 'Invalid file format'}), 400
    
    try:
        workbook = load_workbook(file)
        sheet = workbook.active
        headers = [cell.value for cell in sheet[1]]
        data = []
        
        for row in sheet.iter_rows(min_row=2):
            row_data = [cell.value for cell in row]
            data.append(row_data)
        
        return jsonify({
            'headers': headers,
            'data': data
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/export-pptx', methods=['POST'])
def export_pptx():
    try:
        data = request.json
        prs = Presentation()
        
        # Set slide size based on ratio
        if data['ratio'] == '16:9':
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        elif data['ratio'] == '16:10':
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(8.333)
        else:  # 4:3
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        
        # Create slides based on Excel data
        elements = data['elements']
        excel_data = data.get('excelData', [])
        
        if excel_data and len(excel_data) > 1:
            # Skip header row
            for row_data in excel_data[1:]:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                for element in elements:
                    add_element_to_slide(slide, element, row_data)
        else:
            # Create single slide with current layout
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            for element in elements:
                add_element_to_slide(slide, element)
        
        # Save presentation to memory
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        return send_file(
            pptx_io,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx'
        )
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def add_element_to_slide(slide, element, row_data=None):
    # Convert pixels to inches (1 inch = 72 points)
    PIXELS_PER_INCH = 72
    
    x = float(element['x']) / PIXELS_PER_INCH
    y = float(element['y']) / PIXELS_PER_INCH
    width = float(element['width']) / PIXELS_PER_INCH
    height = float(element['height']) / PIXELS_PER_INCH
    
    if element['type'] == 'image':
        # Handle image elements
        image_data = element['imageData'].split(',')[1]  # Remove data URL prefix
        image_bytes = base64.b64decode(image_data)
        image_stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(image_stream, Inches(x), Inches(y), Inches(width), Inches(height))
    
    else:
        # Handle text elements
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        text_frame = text_box.text_frame
        
        # Get text content
        if 'columnIndex' in element and row_data:
            text = str(row_data[element['columnIndex']])
        else:
            text = element['text']
        
        # Add and format text
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        
        # Apply text formatting
        font = run.font
        font.name = element.get('fontFamily', 'Arial')
        font.size = Pt(float(element.get('fontSize', 12)))
        
        if element.get('color'):
            color = element['color'].lstrip('#')
            font.color.rgb = RGBColor(
                int(color[:2], 16),
                int(color[2:4], 16),
                int(color[4:], 16)
            )
        
        # Apply text alignment
        if element.get('textAlign'):
            if element['textAlign'] == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif element['textAlign'] == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
        
        # Apply text styles
        font.bold = element.get('bold', False)
        font.italic = element.get('italic', False)
        if element.get('underline', False):
            font.underline = True

if __name__ == '__main__':
    app.run(debug=True)